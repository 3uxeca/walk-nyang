"""walk3d v2 업데이트 리포트 PPTX 생성 스크립트.

실행: `python3 docs/build_pptx.py`
산출물: docs/walk3d-v2.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = Path(__file__).resolve().parent
ASSETS = HERE / "assets"
OUT = HERE / "walk3d-v2.pptx"

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Palette
PINK = RGBColor(0xF0, 0x62, 0x92)
PINK_LIGHT = RGBColor(0xFF, 0xB3, 0xC6)
PURPLE = RGBColor(0x6A, 0x4C, 0x93)
BG = RGBColor(0xFE, 0xFA, 0xF3)
TEXT = RGBColor(0x30, 0x30, 0x30)
SUB = RGBColor(0x70, 0x70, 0x70)
ACCENT_BG = RGBColor(0xFF, 0xE0, 0xEC)
ACCENT_TXT = RGBColor(0xC6, 0x28, 0x28)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]


def add_slide():
    slide = prs.slides.add_slide(BLANK)
    # 배경
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    return slide


def add_textbox(slide, left, top, width, height, text, *, size=18, bold=False,
                color=TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Pretendard"
    return tb


def add_header(slide, title, subtitle=None):
    # 좌측 핑크 바
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.5), Inches(0.45),
                                 Inches(0.08), Inches(0.55))
    bar.line.fill.background()
    bar.fill.solid(); bar.fill.fore_color.rgb = PINK
    # 타이틀
    add_textbox(slide, Inches(0.7), Inches(0.4), Inches(11.5), Inches(0.6),
                title, size=28, bold=True, color=PINK)
    if subtitle:
        add_textbox(slide, Inches(0.7), Inches(0.95), Inches(11.5), Inches(0.4),
                    subtitle, size=14, color=SUB)
    # 하단 페이지네이션 받침선
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.5), Inches(7.05),
                                  Inches(12.333), Inches(0.02))
    line.line.fill.background()
    line.fill.solid(); line.fill.fore_color.rgb = PINK_LIGHT


def add_page_number(slide, n, total):
    add_textbox(slide, Inches(11.5), Inches(7.1), Inches(1.3), Inches(0.3),
                f"{n} / {total}", size=10, color=SUB, align=PP_ALIGN.RIGHT)
    add_textbox(slide, Inches(0.5), Inches(7.1), Inches(6), Inches(0.3),
                "walk3d v2 · 업데이트 리포트", size=10, color=SUB)


def add_bullets(slide, left, top, width, height, bullets, *, size=16,
                color=TEXT, indent_px=0):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, is_sub = item
        else:
            text, is_sub = item, False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 1 if is_sub else 0
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = ("◦ " if is_sub else "• ") + text
        run.font.size = Pt(size - (2 if is_sub else 0))
        run.font.color.rgb = color if not is_sub else SUB
        run.font.name = "Pretendard"
    return tb


def add_table(slide, left, top, width, height, rows, col_widths=None,
              header_size=14, body_size=12):
    n_rows = len(rows)
    n_cols = len(rows[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    t = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            t.columns[i].width = w
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = t.cell(r, c)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            if r == 0:
                run.font.bold = True
                run.font.size = Pt(header_size)
                run.font.color.rgb = PURPLE
                cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT_BG
            else:
                run.font.size = Pt(body_size)
                run.font.color.rgb = TEXT
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if r % 2 == 1 else BG
            run.font.name = "Pretendard"
    return table_shape


def add_image(slide, path, left, top, width=None, height=None):
    if not path.exists():
        return None
    kwargs = {}
    if width: kwargs["width"] = width
    if height: kwargs["height"] = height
    return slide.shapes.add_picture(str(path), left, top, **kwargs)


def add_rounded_badge(slide, left, top, width, height, text, *, fill=ACCENT_BG, color=ACCENT_TXT, size=12):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.line.fill.background()
    s.fill.solid(); s.fill.fore_color.rgb = fill
    tf = s.text_frame
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = color
    run.font.name = "Pretendard"
    return s


# ──────────────────────────────────────────────────────────────────────────
# 슬라이드들
# ──────────────────────────────────────────────────────────────────────────

slides_meta = []  # (slide, title)

# Slide 1 — 타이틀
s = add_slide()
# 배경 핑크 블록
hero = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(2.4))
hero.line.fill.background(); hero.fill.solid(); hero.fill.fore_color.rgb = PINK_LIGHT

add_textbox(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.5),
            "walk3d · v2 UPDATE REPORT",
            size=14, bold=True, color=PURPLE)
add_textbox(s, Inches(0.8), Inches(1.0), Inches(11), Inches(1.2),
            "🐈 고양이와 걷는 3D 마을",
            size=44, bold=True, color=PINK)
add_textbox(s, Inches(0.8), Inches(2.1), Inches(11), Inches(0.5),
            "에셋 · UI · 월드 시스템 전면 교체",
            size=18, color=TEXT)

# 중앙 스크린샷
add_image(s, ASSETS / "screen-hearts.png", Inches(3.5), Inches(2.9), height=Inches(3.8))

add_textbox(s, Inches(0.8), Inches(6.7), Inches(11), Inches(0.4),
            "TypeScript · Three.js · Vite · Web Audio     |     2026.04",
            size=12, color=SUB)
slides_meta.append((s, None))


# Slide 2 — 게임 개요
s = add_slide()
add_header(s, "게임 개요",
           "플레이어가 고양이가 되어 절차적으로 생성되는 3D 마을을 걷는 힐링 게임")

add_image(s, ASSETS / "screen-gameplay.png", Inches(7.3), Inches(1.6), height=Inches(5.0))

add_bullets(s, Inches(0.7), Inches(1.7), Inches(6.3), Inches(5.2), [
    "자유 탐험 — 무한하게 이어지는 절차적 청크 월드",
    ("결정적 생성 (same seed → same world)", True),
    ("청크당 건물 · 소품 · 아이템 시드 기반 배치", True),
    "4개 지역 (Meadow / Harbor / Forest / Wildlands)",
    ("레벨 달성 시 자동 언락, 도로로만 이동 허용", True),
    "아이템 3종 (별 · 코인 · 젬) 수집 게임플레이",
    "idle 5초 → 고양이 그르렁 + 춤 + 하트 파티클",
    "브라우저 하나로 플레이 (빌드 ≈ 177KB gzip)",
], size=15)
slides_meta.append((s, "개요"))


# Slide 3 — v1 → v2
s = add_slide()
add_header(s, "v1 → v2 한눈에", "에셋·비주얼·시스템 전방위 교체")

rows = [
    ["영역", "v1", "v2"],
    ["월드 에셋", "전부 절차적 박스", "Kenney City Kit GLB + 절차적 폴백"],
    ["도로", "폭 3유닛 단색 2장", "총 폭 14유닛 · 인도/중앙선/엣지 · non-overlap"],
    ["지형", "단색 MeshToon", "ShaderMaterial 격자 + 지역별 팔레트"],
    ["하늘", "단색 배경", "그라디언트 하늘 돔 + 구름 퍼프 12개"],
    ["사운드", "synth만", "Freesound CC0 meow/purr/footstep 3종 추가"],
    ["UI", "기본 HUD", "컨트롤 HUD · 언락 FX · 수집 이펙트"],
    ["이펙트", "없음", "하트 파티클 · 건물 투명화 · CollectFX"],
    ["세이브", "v1 flat", "schema v2 · 아이템 좌표 버전관리"],
]
add_table(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(5.2), rows,
          col_widths=[Inches(2.0), Inches(3.5), Inches(6.6)],
          header_size=14, body_size=12)
slides_meta.append((s, "v1 vs v2"))


# Slide 4 — 기술 스택
s = add_slide()
add_header(s, "기술 스택", "웹 3D · 시드 기반 생성 · 테스트 커버리지 65개")

# 좌측 런타임
add_textbox(s, Inches(0.7), Inches(1.6), Inches(6), Inches(0.5),
            "런타임", size=20, bold=True, color=PURPLE)
add_bullets(s, Inches(0.7), Inches(2.1), Inches(6), Inches(4.5), [
    "TypeScript 5.x · strict mode",
    "Three.js — WebGLRenderer + PCF Soft Shadow",
    ("GLTFLoader / SkeletonUtils", True),
    ("MeshToonMaterial · ShaderMaterial · Raycaster", True),
    "Web Audio API",
    ("PCM 재생 + 실시간 합성 혼합", True),
    "Web Worker — chunkWorker.ts 청크 프리패스",
], size=14)

# 우측 빌드/품질
add_textbox(s, Inches(7.0), Inches(1.6), Inches(6), Inches(0.5),
            "빌드 / 품질", size=20, bold=True, color=PURPLE)
add_bullets(s, Inches(7.0), Inches(2.1), Inches(6), Inches(4.5), [
    "Vite 6 — dev HMR + prod build",
    "Vitest — 65 test · 9 files 전부 통과",
    "tsc --noEmit 타입 체크",
    "polygonOffset · userData.occludable 등 렌더 테크닉",
    "결정적 RNG (seed-based chunkSeed)",
    "localStorage 세이브 + schema versioning",
], size=14)

# 하단 badges
badges = ["TypeScript", "Three.js", "Vite 6", "Vitest", "Web Audio", "WebGL 2"]
for i, b in enumerate(badges):
    add_rounded_badge(s, Inches(0.7 + i*2.0), Inches(6.4), Inches(1.8), Inches(0.4), b)
slides_meta.append((s, "Tech Stack"))


# Slide 5 — 3D 에셋
s = add_slide()
add_header(s, "사용 에셋 — 3D 모델",
           "CC0 · Kenney / Quaternius 추정 · 전 에셋 머티리얼 깊은 복제 적용")

rows = [
    ["분류", "파일", "개수", "출처 추정"],
    ["캐릭터", "cat/Cat.glb", "1 (스킨드)", "Quaternius RPG Animals"],
    ["건물 (주거)", "suburban/house · apartment.glb", "2", "Kenney City Kit (Suburban)"],
    ["건물 (상업)", "commercial/shop · cafe · tower.glb", "3", "Kenney City Kit (Commercial)"],
    ["소품", "tree · flower · lamp · bench · mailbox.glb", "5", "Kenney Nature / City Kit"],
    ["도로 (미사용)", "roads/road_*.glb", "3", "미사용 (RoadGrid.ts로 절차적 생성)"],
    ["아이템", "(파일 없음)", "0", "코드 내 절차적 생성 (ItemSystem.ts)"],
]
add_table(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(3.6), rows,
          col_widths=[Inches(2.2), Inches(4.2), Inches(1.7), Inches(4.0)])

add_textbox(s, Inches(0.7), Inches(5.5), Inches(12), Inches(0.4),
            "라이선스", size=16, bold=True, color=PURPLE)
add_bullets(s, Inches(0.7), Inches(5.9), Inches(12), Inches(1.3), [
    "Kenney 에셋: CC0 (Public Domain) — https://kenney.nl/assets",
    "Quaternius 에셋(추정): CC0 — https://quaternius.com",
    "정확한 팩 ID/URL은 공개 배포 시 README에 명시 권장 (ASSETS.md 참조)",
], size=13)
slides_meta.append((s, "3D Assets"))


# Slide 6 — 사운드 에셋
s = add_slide()
add_header(s, "사용 에셋 — 사운드", "Freesound.org · Creative Commons 0")

rows = [
    ["파일", "포맷", "용도", "처리"],
    ["meow.wav", "WAV", "아이템 수집 시 야옹", "피치 랜덤 ±10%"],
    ["purring.ogg", "OGG", "idle 5초 그르렁 (loop)", "fade-in 1.5s · fade-out 0.6s"],
    ["footstep.ogg", "OGG", "걷기/대시 발걸음", "피치 랜덤 · 하이패스 200Hz"],
]
add_table(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(2.2), rows,
          col_widths=[Inches(2.3), Inches(1.3), Inches(4.5), Inches(4.0)])

add_textbox(s, Inches(0.7), Inches(4.1), Inches(12), Inches(0.4),
            "Web Audio 합성 (파일 없음)", size=16, bold=True, color=PURPLE)
add_bullets(s, Inches(0.7), Inches(4.5), Inches(12), Inches(2.3), [
    "playJump() — 또잉! 디튠된 사인파 2개 (0.22s)",
    "playDashWhoosh() — sawtooth + 22Hz 비브라토 LFO + 포먼트 필터로 고양이 트릴",
    "playMeow() 폴백 — 2단 bandpass + 사인 비브라토 (파일 로드 실패 시)",
    "playFootstep() 폴백 — 짧은 사인 톤 + exponential decay",
], size=14)
slides_meta.append((s, "Sound Assets"))


# Slide 7 — 월드 시스템
s = add_slide()
add_header(s, "주요 업데이트 ① 월드 시스템",
           "청크 좌표계 근본 수정 + 도로 non-overlap 타일 재설계")

add_textbox(s, Inches(0.7), Inches(1.6), Inches(6), Inches(0.5),
            "청크 좌표계 버그 근본 수정", size=17, bold=True, color=PURPLE)
add_bullets(s, Inches(0.7), Inches(2.1), Inches(6), Inches(2.0), [
    "문제: 아이템이 잠긴 지역에 스폰",
    ("원인: cx·32 중심으로 ±range → 절반이 인접 청크로", True),
    ("해결: (cx+0.5)·32 (청크 내부 중심)", True),
    "세이브 schema v2로 기존 좌표 invalidate",
], size=13)

add_textbox(s, Inches(0.7), Inches(4.4), Inches(6), Inches(0.5),
            "도로 레이아웃 재작성", size=17, bold=True, color=PURPLE)
add_bullets(s, Inches(0.7), Inches(4.9), Inches(6), Inches(2.0), [
    "EW 풀폭 1장 + NS 2토막 (길이 11) — 교차점 겹침 0",
    "EW 인도 2토막 + NS 인도 2토막 (모서리 포함)",
    "엣지/대시 Y 레이어 분리 + polygonOffset",
], size=13)

# 우측 이미지 (town shot — shows roads/sidewalks)
add_image(s, ASSETS / "screen-town.png", Inches(7.2), Inches(1.6), height=Inches(5.2))
slides_meta.append((s, "월드 시스템"))


# Slide 8 — 비주얼 개선
s = add_slide()
add_header(s, "주요 업데이트 ② 비주얼 개선",
           "하늘 · 지형 · Z-fighting · 건물 투명화")

# 좌측 텍스트
add_textbox(s, Inches(0.7), Inches(1.6), Inches(6), Inches(0.4),
            "하늘 / 배경", size=15, bold=True, color=PURPLE)
add_bullets(s, Inches(0.7), Inches(2.0), Inches(6), Inches(1.4), [
    "SkySystem — 반경 180 BackSide 돔",
    "GLSL 그라디언트 셰이더 (topColor · horizonColor)",
    "느리게 오실레이션하는 구름 퍼프 12개",
    "scene.fog 40–110 유닛",
], size=12)

add_textbox(s, Inches(0.7), Inches(3.5), Inches(6), Inches(0.4),
            "지형 / 셰이더", size=15, bold=True, color=PURPLE)
add_bullets(s, Inches(0.7), Inches(3.9), Inches(6), Inches(1.4), [
    "ShaderMaterial 격자 패턴 (UV mod step)",
    "지역별 팔레트 4종 — 언락 색상 연계",
], size=12)

add_textbox(s, Inches(0.7), Inches(5.1), Inches(6), Inches(0.4),
            "Z-fighting 전면 해결", size=15, bold=True, color=PURPLE)
add_bullets(s, Inches(0.7), Inches(5.5), Inches(6), Inches(1.5), [
    "글로벌 300×300 바닥 제거 → 청크별 단일화",
    "도로/인도/엣지/대시 Y 레이어 (0.02/0.03/0.04)",
    "polygonOffset 모든 도로 재질 적용",
], size=12)

# 우측 이미지 (hearts shot)
add_image(s, ASSETS / "screen-hearts.png", Inches(7.2), Inches(1.6), height=Inches(5.2))
slides_meta.append((s, "비주얼 개선"))


# Slide 9 — 게임플레이
s = add_slide()
add_header(s, "주요 업데이트 ③ 게임플레이",
           "물리 · 지역 언락 · 이동 컨트롤")

add_textbox(s, Inches(0.7), Inches(1.6), Inches(12), Inches(0.5),
            "물리 · 충돌", size=17, bold=True, color=PURPLE)
add_bullets(s, Inches(0.7), Inches(2.1), Inches(12), Inches(2.0), [
    "BuildingColliders — 원통형 충돌 레지스트리 (반경 3.8)",
    ("청크 로드/언로드 시 등록/해제 자동", True),
    "X/Z 축 분리 슬라이딩 충돌로 벽 따라 미끄러짐",
    "지형 STEP_HEIGHT = 0.35 이하 자동 오름",
], size=13)

add_textbox(s, Inches(0.7), Inches(4.0), Inches(12), Inches(0.5),
            "지역 언락", size=17, bold=True, color=PURPLE)
add_bullets(s, Inches(0.7), Inches(4.5), Inches(12), Inches(1.3), [
    "4개 지역 · 레벨업으로 자동 언락 + RegionUnlockFX 오버레이",
    "잠긴 영역은 도로 밴드(ROAD_HALF=9)만 진입 허용",
], size=13)

add_textbox(s, Inches(0.7), Inches(5.6), Inches(12), Inches(0.5),
            "대시 · 점프", size=17, bold=True, color=PURPLE)
add_bullets(s, Inches(0.7), Inches(6.1), Inches(12), Inches(1.2), [
    "Shift 대시 8→18 + dashWhoosh 트릴 사운드",
    "Space 점프 JUMP_FORCE 11 · GRAVITY -25 · 착지 감지",
], size=13)
slides_meta.append((s, "게임플레이"))


# Slide 10 — UI / FX
s = add_slide()
add_header(s, "주요 업데이트 ④ UI / FX",
           "HUD · 하트 파티클 · 건물 투명화")

add_image(s, ASSETS / "screen-occlusion.png", Inches(7.4), Inches(1.6), height=Inches(5.2))

add_textbox(s, Inches(0.7), Inches(1.6), Inches(6.5), Inches(0.5),
            "컨트롤 HUD", size=15, bold=True, color=PURPLE)
add_bullets(s, Inches(0.7), Inches(2.0), Inches(6.5), Inches(1.4), [
    "좌상단 수집 진행도 + 현재 지역",
    "하단 WASD / DASH / JUMP 키가이드",
    "키 입력 시 하이라이트 애니메이션",
], size=12)

add_textbox(s, Inches(0.7), Inches(3.5), Inches(6.5), Inches(0.5),
            "HeartFX", size=15, bold=True, color=PURPLE)
add_bullets(s, Inches(0.7), Inches(3.9), Inches(6.5), Inches(1.4), [
    "idle 5초 후 하트 위로 천천히 떠오름",
    "4색 랜덤 · bezier heart · billboard lookAt",
    "fade in 0–20% · hold · fade out 60–100%",
], size=12)

add_textbox(s, Inches(0.7), Inches(5.4), Inches(6.5), Inches(0.5),
            "건물 투명화", size=15, bold=True, color=PURPLE)
add_bullets(s, Inches(0.7), Inches(5.8), Inches(6.5), Inches(1.4), [
    "카메라→고양이 Raycaster",
    "가리는 건물 전체 메시 0.18 opacity",
    "인스턴스별 머티리얼 복제 (전염 버그 해결)",
], size=12)
slides_meta.append((s, "UI / FX"))


# Slide 11 — 버그 수정
s = add_slide()
add_header(s, "해결된 버그 — 3회차 12건", "v2에서 전부 fix · FIXES.md 참조")

rows = [
    ["#", "영역", "이슈", "핵심 해결"],
    ["1", "월드", "도로에 건물이 막힘", "셀 스킵 && → ||"],
    ["2", "에셋", "벤치 너무 작음", "GLTF scale 0.9 → 2.2"],
    ["3", "물리", "고양이가 건물 관통", "원통 충돌 + 슬라이딩"],
    ["4", "월드", "경계 아이템 접근 불가(1차)", "스폰 반경 0.8 → 0.55"],
    ["5", "에셋", "가로등 작고 비대칭", "scale 1.8 → 5.0, 대칭형"],
    ["6", "에셋", "코인이 눕혀서 회전", "rotation.x = π/2 + 토러스 제거"],
    ["7", "에셋", "나무 너무 작음", "scale 2.5 → 4.0"],
    ["8", "렌더", "고양이 시야 방해", "카메라→고양이 레이캐스트 투명화"],
    ["9", "월드", "가로등이 도로에 스폰", "도로 밴드 ±7 유닛 스킵"],
    ["10", "월드", "잠긴 지역 아이템 (근본)", "중심 (cx+0.5)·CHUNK_SIZE"],
    ["11", "비주얼", "도로/지형 단순", "RoadGrid 재작성 + 지형 셰이더"],
    ["12", "렌더", "Z-fighting", "Y 레이어 + polygonOffset"],
]
add_table(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5), rows,
          col_widths=[Inches(0.5), Inches(1.0), Inches(4.5), Inches(6.3)],
          header_size=13, body_size=11)
slides_meta.append((s, "Bug Fixes"))


# Slide 12 — v2 후반 추가
s = add_slide()
add_header(s, "v2 후반 추가 개선",
           "머티리얼 공유 · 바닥 깜빡임 · 도로 교차 겹침")

add_textbox(s, Inches(0.7), Inches(1.6), Inches(12), Inches(0.5),
            "🧵 빌딩 머티리얼 공유 버그",
            size=17, bold=True, color=PINK)
add_bullets(s, Inches(0.7), Inches(2.1), Inches(12), Inches(1.2), [
    "증상: 고양이 앞 건물 외에 멀리 있는 동일 타입 건물도 함께 투명해짐",
    "원인: THREE.Group.clone()은 머티리얼을 공유 — 한 인스턴스의 opacity 변경이 전파",
    "해결: AssetManager.clone()에서 머티리얼 깊은 복제 (material.clone())",
], size=12)

add_textbox(s, Inches(0.7), Inches(3.7), Inches(12), Inches(0.5),
            "🟢 바닥 전체 깜빡임 (심각)",
            size=17, bold=True, color=PINK)
add_bullets(s, Inches(0.7), Inches(4.2), Inches(12), Inches(1.2), [
    "main.ts의 300×300 글로벌 바닥과 청크별 32×32 바닥이 둘 다 y=0.0에 coplanar",
    "이동 중 카메라 변경 → 전 영역 z-fighting",
    "해결: 글로벌 바닥 제거 → 청크 바닥으로 단일화",
], size=12)

add_textbox(s, Inches(0.7), Inches(5.7), Inches(12), Inches(0.5),
            "🛣️ 도로 EW/NS 교차점 겹침",
            size=17, bold=True, color=PINK)
add_bullets(s, Inches(0.7), Inches(6.2), Inches(12), Inches(1.0), [
    "NS 도로를 EW 도로를 피해 두 토막(길이 11)으로 분할 — 10×10 겹침 제거",
], size=12)
slides_meta.append((s, "Late Fixes"))


# Slide 13 — 아키텍처
s = add_slide()
add_header(s, "아키텍처 요약", "시드 기반 결정적 월드 · 모듈형 시스템")

arch = """main.ts (게임 루프)
├─ Character · Controller · ThirdPersonCamera
├─ ChunkManager → ChunkGenerator → ChunkMeshFactory
│                                   ├─ RoadGrid
│                                   ├─ Buildings ┐
│                                   ├─ Props     ├ AssetManager (GLB + 폴백)
│                                   └─ Items     ┘
├─ Terrain · BuildingColliders
├─ ItemSystem · ProgressSystem · RegionManager
├─ SaveSystem (localStorage + schema v2)
├─ SoundSystem (Web Audio)
├─ HUD · ControlsHUD · RegionUnlockFX
├─ CollectFX · HeartFX · SkySystem
└─ WORLD_SEED (결정적 월드)"""

tb = s.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.8))
tb.fill.solid(); tb.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
tb.line.color.rgb = PINK_LIGHT
tf = tb.text_frame
tf.word_wrap = True
tf.margin_top = Inches(0.15); tf.margin_bottom = Inches(0.15)
tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.25)
for i, line in enumerate(arch.split("\n")):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    run = p.add_run()
    run.text = line
    run.font.name = "Menlo"
    run.font.size = Pt(13)
    run.font.color.rgb = TEXT

add_textbox(s, Inches(0.8), Inches(6.6), Inches(12), Inches(0.4),
            "결정적 생성 — 같은 seed는 항상 같은 월드 (플레이어 간 공유 가능)",
            size=13, color=PURPLE, bold=True)
slides_meta.append((s, "Architecture"))


# Slide — OMC 파이프라인
s = add_slide()
add_header(s, "개발 프로세스 — OMC 파이프라인",
           "Deep Interview → RALPLAN-DR Consensus → Ralph PRD Loop")

# 왼쪽: 파이프라인 다이어그램
add_textbox(s, Inches(0.7), Inches(1.6), Inches(6), Inches(0.5),
            "3단계 자동화 파이프라인", size=17, bold=True, color=PURPLE)

stages = [
    ("1. Deep Interview",   "요구사항 6라운드 인터뷰",  "모호도 18% PASSED",    PINK),
    ("2. RALPLAN-DR",       "Planner → Architect → Critic", "합의 기반 계획 수립",  PURPLE),
    ("3. Ralph PRD Loop",   "12개 스토리 순차 실행",    "자동 검증 + 자동 재시도", PINK),
]
for i, (title, desc, result, color) in enumerate(stages):
    top = Inches(2.2 + i * 1.45)
    # 번호 박스
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(0.7), top, Inches(6), Inches(1.25))
    box.line.fill.background()
    box.fill.solid()
    box.fill.fore_color.rgb = ACCENT_BG if i % 2 == 0 else RGBColor(0xF5, 0xF0, 0xFF)
    add_textbox(s, Inches(0.9), top + Inches(0.1), Inches(5.7), Inches(0.4),
                title, size=14, bold=True, color=color)
    add_textbox(s, Inches(0.9), top + Inches(0.5), Inches(5.7), Inches(0.3),
                desc, size=11, color=TEXT)
    add_textbox(s, Inches(0.9), top + Inches(0.85), Inches(5.7), Inches(0.3),
                "→ " + result, size=11, color=SUB)

# 오른쪽: 실제 산출물 트리
add_textbox(s, Inches(7.2), Inches(1.6), Inches(6), Inches(0.5),
            "실제 산출물 (`.omc/`)", size=17, bold=True, color=PURPLE)

tb = s.shapes.add_textbox(Inches(7.2), Inches(2.1), Inches(5.8), Inches(4.8))
tb.fill.solid(); tb.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
tb.line.color.rgb = PINK_LIGHT
tf = tb.text_frame
tf.word_wrap = True
tf.margin_top = Inches(0.2); tf.margin_bottom = Inches(0.2)
tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2)
tree = """.omc/
├─ specs/
│    └─ deep-interview-walk3d.md
│       (6 rounds · ambiguity 18%)
├─ plans/
│    ├─ walk3d-village-game.md
│    │   (RALPLAN-DR consensus)
│    └─ open-questions.md
├─ prd.json
│    (12 stories · all passes:true)
├─ progress.txt
│    (ralph iteration log)
└─ state/, sessions/
     (OMC runtime state)"""
for i, line in enumerate(tree.split("\n")):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    run = p.add_run()
    run.text = line
    run.font.name = "Menlo"
    run.font.size = Pt(11)
    run.font.color.rgb = TEXT
slides_meta.append((s, "OMC Pipeline"))


# Slide — 멀티에이전트
s = add_slide()
add_header(s, "멀티에이전트 협업",
           "단일 LLM이 아닌 역할별 전문 에이전트 조합 사용")

rows = [
    ["에이전트", "역할", "사용 단계", "모델 티어"],
    ["planner", "요구사항 → 계획 수립 · 12개 스토리 분해", "Deep Interview · RALPLAN", "opus"],
    ["architect", "아키텍처/경계 검토 · steelman 반대", "RALPLAN consensus loop", "opus"],
    ["critic", "원칙/옵션/리스크/검증 일관성 검증", "RALPLAN consensus loop", "opus"],
    ["executor", "스토리별 실제 구현 (코드 작성)", "Ralph 반복 루프", "sonnet / opus"],
    ["verifier", "기수용 기준(AC) 충족 여부 검증", "Ralph post-impl", "sonnet"],
    ["code-reviewer", "스타일·SOLID·회귀 리스크", "Ralph approval pass", "opus"],
    ["ai-slop-cleaner", "구현 후 반복/데드코드 정리", "Ralph mandatory deslop", "skill"],
]
add_table(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(4.5), rows,
          col_widths=[Inches(2.0), Inches(5.3), Inches(3.5), Inches(1.5)],
          header_size=13, body_size=11)

add_textbox(s, Inches(0.7), Inches(6.3), Inches(12), Inches(0.4),
            "작성/리뷰 분리 — 같은 컨텍스트에서 자기 승인 금지. writer-pass와 reviewer-pass를 다른 에이전트로 분리",
            size=11, color=SUB, bold=True)
slides_meta.append((s, "Multi-Agent"))


# Slide — 학습된 프로젝트 룰
s = add_slide()
add_header(s, "프로젝트 메모리 · 룰 / 스킬",
           "작업 중 축적된 규칙 · 이후 세션에 자동 적용")

add_textbox(s, Inches(0.7), Inches(1.6), Inches(6), Inches(0.5),
            "Project Memory (.omc/project-memory.json)", size=15, bold=True, color=PURPLE)
add_bullets(s, Inches(0.7), Inches(2.1), Inches(6), Inches(2.5), [
    "Hot paths — main.ts (47×), Character.ts (28×), SoundSystem.ts (25×)",
    "JS/TS · vitest · npm · build=`npm run build`",
    "세션 간 지속 — 코드 컨벤션/핫패스 자동 recall",
], size=12)

add_textbox(s, Inches(0.7), Inches(4.7), Inches(6), Inches(0.5),
            "Feedback 메모리 (user preferences)", size=15, bold=True, color=PURPLE)
add_bullets(s, Inches(0.7), Inches(5.2), Inches(6), Inches(2.0), [
    "작업 끝나면 한글로 설명",
    "커밋은 논리 단위별로 쪼개서 진행",
    "스크린샷/비디오 검증 후 완료 선언",
], size=12)

add_textbox(s, Inches(7.0), Inches(1.6), Inches(6), Inches(0.5),
            "사용한 OMC Skills", size=15, bold=True, color=PURPLE)
add_bullets(s, Inches(7.0), Inches(2.1), Inches(6), Inches(5), [
    "plan --consensus --di — 합의 기반 계획",
    ("Z-fighting 수정 · 비주얼 개선 계획", True),
    "ralph — PRD 12 스토리 영구 루프",
    ("US-001 ~ US-012 전부 passes:true", True),
    "deep-interview — 요구사항 모호도 게이팅",
    ("18% PASSED (threshold 20%)", True),
    "ai-slop-cleaner — 구현 후 정리",
    "visual-verdict — 스크린샷 QA",
    "auto memory — 사용자 선호/핫패스 기록",
], size=12)
slides_meta.append((s, "Rules & Skills"))


# Slide 14 — 테스트
s = add_slide()
add_header(s, "테스트 · 검증", "65 tests · 9 files · 100% pass")

rows = [
    ["테스트 파일", "케이스 수", "커버 영역"],
    ["SaveSystem.test.ts", "7", "세이브 로드 · 버전 마이그레이션"],
    ["RegionManager.test.ts", "15", "지역 언락 조건"],
    ["ProgressSystem.test.ts", "12", "레벨 · 임계값"],
    ["ItemSystem.test.ts", "9", "아이템 수집 판정"],
    ["ChunkGenerator.test.ts", "4", "결정적 생성"],
    ["noise.test.ts", "6", "시드 노이즈 안정성"],
    ["pool.test.ts", "5", "오브젝트 풀링"],
    ["chunkWorker.test.ts", "2", "워커 메시지"],
    ["Controller.test.ts", "5", "입력 매핑"],
    ["합계", "65", "9 files all pass"],
]
add_table(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(5.0), rows,
          col_widths=[Inches(3.5), Inches(1.5), Inches(7.1)],
          header_size=14, body_size=12)

add_textbox(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.3),
            "npm run build — TypeScript OK · Vite prod build 672KB → 177KB (gzip)",
            size=11, color=SUB)
slides_meta.append((s, "Testing"))


# Slide 15 — 로드맵
s = add_slide()
add_header(s, "로드맵 · 다음 단계", "v3 방향")

add_textbox(s, Inches(0.7), Inches(1.6), Inches(4), Inches(0.5),
            "플레이 경험", size=16, bold=True, color=PINK)
add_bullets(s, Inches(0.7), Inches(2.1), Inches(4), Inches(4.5), [
    "일기/퀘스트 시스템",
    ("마을 NPC 고양이와 대화", True),
    "날씨 · 시간",
    ("밤/낮 셰이더 · 비/눈", True),
    "의상/액세서리",
    ("세이브 연계 커스터마이즈", True),
], size=12)

add_textbox(s, Inches(4.9), Inches(1.6), Inches(4), Inches(0.5),
            "기술", size=16, bold=True, color=PINK)
add_bullets(s, Inches(4.9), Inches(2.1), Inches(4), Inches(4.5), [
    "코드 스플리팅",
    ("dynamic import 청크화", True),
    "압축 텍스처",
    ("KTX2 / Basis 적용", True),
    "모바일 컨트롤",
    ("가상 조이스틱 · 터치", True),
], size=12)

add_textbox(s, Inches(9.1), Inches(1.6), Inches(4), Inches(0.5),
            "콘텐츠", size=16, bold=True, color=PINK)
add_bullets(s, Inches(9.1), Inches(2.1), Inches(4), Inches(4.5), [
    "지역 5–8개 확장",
    ("랜드마크 · 히든 장소", True),
    "지역별 팔레트/건물",
    ("지역 아이덴티티 강화", True),
    "계절 BGM",
    ("지역별 오프셋 루프", True),
], size=12)
slides_meta.append((s, "Roadmap"))


# Slide 16 — 엔딩
s = add_slide()
hero = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
hero.line.fill.background(); hero.fill.solid(); hero.fill.fore_color.rgb = PINK_LIGHT

add_textbox(s, Inches(0), Inches(1.5), SLIDE_W, Inches(1.2),
            "감사합니다 🐾", size=60, bold=True, color=PINK,
            align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0), Inches(3.0), SLIDE_W, Inches(0.5),
            "가만히 멈춰서 고양이를 지켜보기만 해도 행복해요",
            size=20, color=PURPLE, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0), Inches(4.2), SLIDE_W, Inches(0.4),
            "GitHub · github.com/3uxeca/walk-nyang",
            size=16, color=TEXT, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0), Inches(4.8), SLIDE_W, Inches(0.4),
            "Tech · TypeScript · Three.js · Vite · Web Audio",
            size=14, color=SUB, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0), Inches(5.3), SLIDE_W, Inches(0.4),
            "Assets · Kenney CC0 · Freesound CC0 · Quaternius CC0",
            size=14, color=SUB, align=PP_ALIGN.CENTER)

# Cat image center bottom
add_image(s, ASSETS / "screen-cat.png", Inches(5.5), Inches(5.8), height=Inches(1.5))
slides_meta.append((s, None))


# ──────────────────────────────────────────────────────────────────────────
# 페이지 번호 (타이틀/엔딩 제외)
# ──────────────────────────────────────────────────────────────────────────
total = len(slides_meta)
for i, (slide, title) in enumerate(slides_meta, start=1):
    if title is None:
        continue
    add_page_number(slide, i, total)

prs.save(str(OUT))
print(f"saved: {OUT}")
print(f"slides: {total}")
